from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import MSO_ANCHOR
from pptx.dml.color import RGBColor

# Finsurge brand colors
NAVY_BLUE = RGBColor(31, 67, 128)      # #1F4380
PURPLE = RGBColor(163, 31, 104)        # #A31F68
ACCENT_PINK = RGBColor(234, 61, 117)   # #EA3D75
WHITE = RGBColor(255, 255, 255)
LIGHT_GRAY = RGBColor(245, 246, 248)   # #F5F6F8
DARK_GRAY = RGBColor(55, 65, 81)       # #374151

def add_logo(slide, logo_path):
    """Add Finsurge logo to top-right"""
    left = Inches(8.8)
    top = Inches(0.3)
    height = Inches(0.6)
    slide.shapes.add_picture(logo_path, left, top, height=height)

def add_title_slide(prs, title, subtitle, logo_path):
    """Professional title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background: Navy blue
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = NAVY_BLUE

    # Add Finsurge logo
    add_logo(slide, logo_path)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.5), Inches(8), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.8), Inches(4.2), Inches(8), Inches(2))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.word_wrap = True
    p = subtitle_frame.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(24)
    p.font.color.rgb = ACCENT_PINK

    # Tagline
    tagline_box = slide.shapes.add_textbox(Inches(0.8), Inches(6.5), Inches(8), Inches(0.8))
    tagline_frame = tagline_box.text_frame
    p = tagline_frame.paragraphs[0]
    p.text = "Enterprise Risk Management System for Indian & Bhutanese Banking"
    p.font.size = Pt(16)
    p.font.italic = True
    p.font.color.rgb = RGBColor(200, 200, 200)

def add_content_slide(prs, title, content_items, logo_path):
    """Professional content slide with McKinsey styling"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # White background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE

    # Top bar (navy blue accent)
    top_bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.08))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = NAVY_BLUE
    top_bar.line.color.rgb = NAVY_BLUE

    # Add logo
    add_logo(slide, logo_path)

    # Title with accent line
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.35), Inches(8.2), Inches(0.95))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    title_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(38)
    p.font.bold = True
    p.font.color.rgb = NAVY_BLUE

    # Accent line under title
    line = slide.shapes.add_connector(1, Inches(0.8), Inches(1.4), Inches(3.5), Inches(1.4))
    line.line.color.rgb = ACCENT_PINK
    line.line.width = Pt(4)

    # Content box
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.75), Inches(8.4), Inches(5.25))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    text_frame.vertical_anchor = MSO_ANCHOR.TOP

    for i, item in enumerate(content_items):
        if i > 0:
            text_frame.add_paragraph()
        p = text_frame.paragraphs[i]

        # Check if it's a sub-item (starts with 4 spaces)
        if item.startswith("    "):
            p.text = item.strip()
            p.font.size = Pt(16)
            p.level = 1
        else:
            p.text = item
            p.font.size = Pt(18)
            p.level = 0

        p.font.color.rgb = DARK_GRAY
        p.space_before = Pt(6)
        p.space_after = Pt(6)

        # Bold first part before dash for sub-items
        if " — " in p.text:
            runs = p.runs
            if runs:
                parts = p.text.split(" — ")
                if len(parts) == 2:
                    p.text = ""
                    run1 = p.add_run()
                    run1.text = parts[0]
                    run1.font.bold = True
                    run2 = p.add_run()
                    run2.text = " — " + parts[1]

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

logo_path = r"d:\Users\lsgop\Downloads\finsurgelogo.png"

# Slide 1: Title
add_title_slide(prs, "FinsurgeENRIMS", "Real-Time Fraud Detection & Risk Management", logo_path)

# Slide 2: System Overview
add_content_slide(prs, "FRIMS — 10 Core Modules at a Glance", [
    "Fraud Risk Intelligence Management System — purpose-built for Bhutanese banking & RMA regulations",
    "Transaction Monitoring — 26+ configurable rules evaluated on every transaction in real time",
    "Alert Management — AI risk scoring (0–100), SLA enforcement, skill-based auto-assignment",
    "Case Management — Customer 360°, evidence repository, maker-checker, court & recovery tracking",
    "SMS OTP Approval — High-value/high-risk transfers held; customer approves via SMS before settlement",
    "CTR / LVTR Reporting — Auto RMA-format PDFs; Nu. 100,000 threshold; 7-day filing deadline",
    "STR / SAR Reporting — Suspicion-to-FIU-Bhutan workflow with maker-checker sign-off",
    "Police FIR — BPC 2004 FIR generation with MLPCA 2018 / BICMA 2018 references",
    "Watchlist Screening — CID/BRN/PEP screening; OFAC, UN Security Council, RMA watchlists",
    "AI Investigation Copilot — Risk explanation, similar case search, next-step recommendations",
    "Audit & Governance — Immutable log, RBAC, PII encryption, data masking",
], logo_path)

# Slide 3: Multi-Regulator Compliance
add_content_slide(prs, "Multi-Regulator Compliance Framework", [
    "Dual compliance scorecard: RBI (India) + RMA (Bhutan) requirements",
    "RBI Compliance (India)",
    "    4 sections covering fraud risk, data governance, reporting, IT security",
    "    16 requirements auto-tracked with live system metrics",
    "RMA Compliance (Bhutan)",
    "    3 sections covering AML/CFT, regulatory reporting, corporate governance",
    "    14 requirements including CDD, EDD, sanctions screening, STR filing",
    "Integrated risk scoring with multi-jurisdictional support",
    "Automatic compliance assessment with coverage metrics"
], logo_path)

# Slide 4: Investigation Copilot
add_content_slide(prs, "Investigation Copilot — AI-Powered Analysis", [
    "AI-assisted investigation platform accelerating alert analysis and case resolution",
    "Risk Score Transparency",
    "    Hover tooltips show exact risk score formula: Base Alert + Risk Factors + PEP Flag + Severity",
    "    Helps investigators understand why alerts are prioritized and which cases need immediate attention",
    "Ask the Copilot — Conversational Investigation",
    "    'Why is this customer flagged?' — Copilot explains detection pattern and rule triggers",
    "    'Show similar cases' — Retrieves historical closed cases with matching characteristics",
    "    'What are next steps?' — Recommends investigation actions based on alert profile",
    "Context Integration: Previous cases, activity timeline, related alerts, network relationships in one unified view"
], logo_path)

# Slide 5: Regulatory Document Generation
add_content_slide(prs, "Regulatory Document Generation", [
    "Multi-jurisdictional compliance document automation with format selection",
    "CTR/LCTR — Cash Transaction Reporting",
    "    RBI Format: FIU-IND reference, INR 10 Lakh threshold, India compliance",
    "    RMA Format: FIU-Bhutan reference, Nu. 100,000 threshold, Bhutan AML/CFT Rules 2009",
    "STR/SAR — Suspicious Transaction Reporting",
    "    RBI Format: PMLA-compliant STR with principal officer workflow",
    "    RMA Format: AML/CFT Rules 2009, Section 11, 7-day filing deadline",
    "LVTR — Large Value Transaction Report (RMA Rule 14)",
    "    Track Nu. 100,000+ transactions; 7-day filing window; full 5-year audit trail",
    "FIR — Police First Information Reports",
    "    India Police (IPC sections); Bhutan Police (BPC 2004 sections with MLPCA/BICMA references)"
], logo_path)

# Slide 6: SMS OTP Transaction Approval
add_content_slide(prs, "SMS OTP — Real-Time Transaction Approval", [
    "Unique CBS-integrated OTP gate: high-value or high-risk transfers are held until customer confirms",
    "Approval Thresholds",
    "    Auto-approve — below Nu. 25,000 (low risk, standard retail transfer)",
    "    SMS OTP required — Nu. 25,001 to Nu. 500,000 (customer must confirm within 5 minutes)",
    "    Analyst pre-approval — above Nu. 500,000 (compliance review before release)",
    "How It Works",
    "    FRIMS signals CBS: HOLD — funds reserved, not settled",
    "    OTP dispatched to registered mobile via Tashi Cell / B-Mobile (SHA-256 stored, 3-attempt limit)",
    "    Customer approves → FRIMS signals CBS: RELEASE — transfer settles normally",
    "    Customer rejects or timeout → FRIMS signals CBS: CANCEL — alert escalated to analyst",
    "Full audit trail: OTP hash, hold/release/cancel timestamps, CBS response codes",
], logo_path)

# Slide 7: Data Integration & Auth Logs
add_content_slide(prs, "Data Integration & Cyber Fraud Detection", [
    "23 system interfaces feeding the FRIMS detection engine — real-time and batch",
    "Core Banking System (REST API / ISO 8583)",
    "    Primary data source: transactions, accounts, balances, beneficiary details",
    "    Real-time event push (preferred) or daily CSV batch with reconciliation",
    "Authentication & Authorisation Logs — Cyber Fraud Detection",
    "    Login attempts, OTP events, password/PIN resets → triggers SIM swap, ATO, phishing rules",
    "    Device registration changes, session anomalies → triggers geo-anomaly, remote-access rules",
    "    Maker-checker audit trail, privileged user activities → triggers internal fraud rules",
    "Other Key Interfaces",
    "    Watchlists: OFAC, UN Security Council, RMA — daily refresh for PEP/sanctions screening",
    "    FIU-Bhutan SFTP Portal · Royal Bhutan Police (RBP) · CID/BRN Registry API · SMS Gateway",
], logo_path)

# Slide 9: RMA AML/CFT Controls
add_content_slide(prs, "RMA AML/CFT Compliance (Bhutan)", [
    "Comprehensive Anti-Money Laundering & Counter-Financing of Terrorism controls",
    "Customer Due Diligence (CDD) & Enhanced Due Diligence (EDD)",
    "    Full KYC verification at account opening with document validation",
    "    Automated PEP screening for politically exposed persons",
    "    Ultimate Beneficial Owner (UBO) identification and verification",
    "Ongoing Transaction Monitoring",
    "    Real-time screening of all transactions against risk patterns",
    "    Sanctions screening against OFAC, UN, and FATF watchlists",
    "Suspicious Transaction Reporting (STR) to FIU-Bhutan within 7 days"
], logo_path)

# Slide 10: RMA Regulatory Reporting
add_content_slide(prs, "RMA Regulatory Reporting & LVTR", [
    "Automated filing with RMA and Financial Intelligence Unit - Bhutan",
    "Suspicious Transaction Reporting (STR)",
    "    Integrated STR workflow with 7-day filing deadline to FIU-Bhutan",
    "    Comprehensive case documentation and investigation tracking",
    "Large Value Transaction Reporting (LVTR)",
    "    Track and report transactions exceeding Nu. 100,000 threshold",
    "    Full transaction records maintained for 5-year RMA requirement",
    "Annual AML/CFT Compliance Report",
    "    Board-level reporting with metrics and incident summary"
], logo_path)

# Slide 11: Compliance Scorecard Dashboard
add_content_slide(prs, "Compliance Scorecard — Live Metrics", [
    "Real-time compliance tracking with auto-calculated scores",
    "RBI Scorecard: 16 requirements across 4 sections",
    "    Fraud Risk Management, Data Governance, Reporting, IT Security",
    "RMA Scorecard: 14 requirements across 3 sections",
    "    AML/CFT, Regulatory Reporting, Corporate Governance",
    "Each requirement shows:",
    "    Status: Compliant / Partial / At-Risk / Not Implemented",
    "    Coverage %: Based on actual system state and metrics",
    "    Evidence: Specific data supporting compliance status"
], logo_path)

# Slide 12: Alert Management & Investigation
add_content_slide(prs, "Investigation & Case Management", [
    "Complete alert lifecycle with SLA enforcement and intelligent routing",
    "Alert States: New → Assigned → Under Review → Escalated → Closed",
    "Severity-Based SLAs: Critical (4h), High (24h), Medium (72h), Low (168h)",
    "Intelligent Auto-Assignment: Routes alerts to available analysts by skill, workload",
    "Investigation Features",
    "    Add detailed notes, attach evidence documents, link related cases",
    "    Bulk actions: Close false positives, escalate groups, reassign batches",
    "Dashboard Analytics: New alerts trending, SLA compliance %, closure patterns"
], logo_path)

# Slide 13: Case Management & Document Downloads
add_content_slide(prs, "Case Management & Document Downloads", [
    "Consolidate related alerts and evidence into formal investigation cases",
    "Document Generation & Download",
    "    Police FIR — Download full investigation report with all case details",
    "    CTR/SAR Filings — Download regulatory reports for archival or reprint",
    "    Generate formatted .txt documents suitable for filing and compliance",
    "Evidence Management: Attach documents, screenshots, transaction exports",
    "Case Types: Fraud, Cyber Attack, Money Laundering, Operational Risk, Internal Misconduct",
    "Case Lifecycle: Open → Assigned → Under Investigation → Escalated → Closed"
], logo_path)

# Slide 14: Regulatory Filing & Dashboard
add_content_slide(prs, "Regulatory Filing & Real-Time Monitoring", [
    "Automated filing workflow with document generation and downloads",
    "CTR/SAR Filings: Download reports with complete filing details",
    "FIR Management: File reports with police, track progress, download documents",
    "Executive visibility into risk management operations",
    "Key Performance Indicators (Real-Time)",
    "    New alerts (24h), Overdue alerts, New cases, Closure rate, SLA compliance %",
    "Operational Dashboards",
    "    Alert heatmap by hour, customer segment, rule category",
    "    Team workload and case resolution metrics"
], logo_path)

# Slide 12: Audit Trail & Security
add_content_slide(prs, "Audit Trail, Security & Governance", [
    "Immutable audit logging for regulatory compliance (RBI + RMA requirements)",
    "Audit Coverage",
    "    Detection: Rule creation, alert generation, escalation triggers",
    "    Investigation: Case creation, evidence, notes, disposition",
    "    Access: User login, data access, rule changes, document downloads",
    "    Retention: 5 years (RBI/RMA mandate) for audit logs, 7 years transaction data",
    "Security Controls",
    "    Role-Based Access with least privilege; JWT-based authentication",
    "    Data Protection: PII encryption, API response masking, HTTPS-only",
    "    Compliance Ready: Meets RBI Master Direction, RMA AML/CFT Rules"
], logo_path)

# Save
out = r"D:/Users/lsgop/Downloads/FinsurgeENRIMS_MultiRegulator_RMA.pptx"
prs.save(out)
print(f"[OK] PowerPoint saved: {out}")
print("[SLIDES] 16 functionality-only slides:")
print("  1  Title")
print("  2  FRIMS — 10 Core Modules")
print("  3  Multi-Regulator Compliance Framework")
print("  4  Investigation Copilot")
print("  5  Regulatory Document Generation (CTR/LVTR/STR/FIR)")
print("  6  SMS OTP Transaction Approval  [NEW]")
print("  7  Data Integration & Cyber Fraud Detection  [NEW]")
print("  8  RMA AML/CFT Compliance")
print("  9  RMA Regulatory Reporting & LVTR")
print(" 10  Compliance Scorecard — Live Metrics")
print(" 11  Investigation & Case Management")
print(" 12  Case Management & Document Downloads")
print(" 13  Regulatory Filing & Real-Time Monitoring")
print(" 14  Audit Trail, Security & Governance")
