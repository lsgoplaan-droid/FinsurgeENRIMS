from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
import os

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Finsurge colors
PRIMARY_BLUE = RGBColor(75, 85, 180)
ACCENT_ORANGE = RGBColor(255, 102, 0)
DARK_TEXT = RGBColor(40, 40, 40)
LIGHT_BG = RGBColor(245, 247, 250)
WHITE = RGBColor(255, 255, 255)

def add_title_slide(title, subtitle, logo_path=None):
    """Add title slide with logo"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE

    # Add gradient-like background with colored bar
    shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = PRIMARY_BLUE
    shape.line.color.rgb = PRIMARY_BLUE

    # Add logo if exists
    if logo_path and os.path.exists(logo_path):
        try:
            slide.shapes.add_picture(logo_path, Inches(0.5), Inches(0.5), width=Inches(1.5))
        except:
            pass

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(54)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = WHITE

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(9), Inches(3))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.word_wrap = True
    subtitle_frame.text = subtitle
    subtitle_frame.paragraphs[0].font.size = Pt(28)
    subtitle_frame.paragraphs[0].font.color.rgb = DARK_TEXT

def add_content_slide(title, sections, logo_path=None):
    """Add content slide with multiple sections"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = LIGHT_BG

    # Add logo if exists
    if logo_path and os.path.exists(logo_path):
        try:
            slide.shapes.add_picture(logo_path, Inches(9), Inches(0.3), width=Inches(0.7))
        except:
            pass

    # Title bar
    title_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.8))
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = PRIMARY_BLUE
    title_shape.line.color.rgb = PRIMARY_BLUE

    # Title text
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(8), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(40)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = WHITE

    # Add sections
    y_position = 1.2
    for i, section in enumerate(sections):
        # Section heading with orange accent
        accent = slide.shapes.add_shape(1, Inches(0.3), Inches(y_position), Inches(0.1), Inches(0.4))
        accent.fill.solid()
        accent.fill.fore_color.rgb = ACCENT_ORANGE
        accent.line.color.rgb = ACCENT_ORANGE

        heading_box = slide.shapes.add_textbox(Inches(0.6), Inches(y_position), Inches(9), Inches(0.35))
        heading_frame = heading_box.text_frame
        heading_frame.text = section['heading']
        heading_frame.paragraphs[0].font.size = Pt(18)
        heading_frame.paragraphs[0].font.bold = True
        heading_frame.paragraphs[0].font.color.rgb = PRIMARY_BLUE

        # Section content
        content_box = slide.shapes.add_textbox(Inches(0.6), Inches(y_position + 0.4), Inches(9), Inches(1.2))
        content_frame = content_box.text_frame
        content_frame.word_wrap = True
        content_frame.text = section['content']
        for para in content_frame.paragraphs:
            para.font.size = Pt(14)
            para.font.color.rgb = DARK_TEXT
            para.space_before = Pt(3)
            para.space_after = Pt(3)

        y_position += 1.6
        if y_position > 6.5:
            break

logo_path = r"D:\Users\lsgop\Downloads\finsurgelogo.png"

# ====== SLIDE 1: TITLE SLIDE ======
add_title_slide(
    "FinsurgeENRIMS",
    "Enterprise Risk Management System\nReal-time Detection, Investigation & Regulatory Compliance for Banking",
    logo_path
)

# ====== SLIDE 2: DASHBOARD ======
add_content_slide(
    "Dashboard — Real-Time Oversight",
    [
        {
            'heading': '📊 Executive KPIs',
            'content': 'Centralized view of critical metrics: open alerts (30), active cases (12), fraud loss (₹53.5M), overdue SLA (2 cases). Color-coded severity indicators for quick identification of risk areas.'
        },
        {
            'heading': '📈 Trend Analysis',
            'content': 'Visualize alert volume, case closure rates, and regulatory filing trends over time. Drill-down capability to trace metrics back to source transactions and customers.'
        },
        {
            'heading': '⚡ Instant Actions',
            'content': 'Click any metric to filter and investigate. Create cases directly from alerts. Assign alerts to analysts with SLA tracking (response time: 2-24 hours based on priority).'
        }
    ],
    logo_path
)

# ====== SLIDE 3: ALERT MANAGEMENT ======
add_content_slide(
    "Alert Detection & Investigation",
    [
        {
            'heading': '🎯 Automated Detection',
            'content': '26 detection rules monitor transactions 24/7 for suspicious patterns: structuring, dormant account activation, unusual velocity, high-risk geography. Real-time risk scoring (0-100) on every transaction.'
        },
        {
            'heading': '🔍 Investigation Copilot',
            'content': 'AI-powered analysis of new alerts. Automatic context gathering: customer profile, account history, network relationships. Smart assignment to available analysts with workload balancing.'
        },
        {
            'heading': '✅ Alert Response',
            'content': 'Mark as true positive → escalate to cases. False positive → close with reason. Escalated → manager review. Full audit trail with user, timestamp, and notes on every status change.'
        }
    ],
    logo_path
)

# ====== SLIDE 4: CASE MANAGEMENT ======
add_content_slide(
    "Case Management & Escalation",
    [
        {
            'heading': '📋 Full Case Lifecycle',
            'content': 'Create case from alert → assign to investigator → track findings → escalate if needed → set disposition (true/false positive) → file regulatory report. SLA tracking with visual indicators for overdue cases.'
        },
        {
            'heading': '🔗 Evidence & Context',
            'content': 'Attach evidence (documents, screenshots, investigator notes). Link related alerts and transactions. Build investigation narrative with automated timeline from case activity log.'
        },
        {
            'heading': '⚖️ Regulatory Workflow',
            'content': 'Cases flagged as true positive auto-file CTR/SAR to FIU-IND portal. Track filing status and regulatory feedback. Maintain 5-year audit trail for compliance reviews.'
        }
    ],
    logo_path
)

# ====== SLIDE 5: CUSTOMER 360 & FRAUD DETECTION ======
add_content_slide(
    "Customer 360 & Fraud Detection",
    [
        {
            'heading': '👤 Customer 360 Profile',
            'content': 'Complete customer view: verification status, risk category, accounts, transaction history (last 50), related alerts, open cases, network relationships. PEP (politically exposed person) flagging with auto-enhanced monitoring for high-risk profiles.'
        },
        {
            'heading': '💳 Fraud Detection',
            'content': 'Card fraud monitoring: velocity checks (amount/frequency), cross-border thresholds, merchant category anomalies. Transaction scoring integrates customer risk + behavioral patterns + rules engine output.'
        },
        {
            'heading': '📊 Risk Scoring',
            'content': 'Dynamic risk adjustment: base risk (low/medium/high/very_high) + behavioral modifiers = real-time customer risk score (0-100). Influences alert sensitivity and SLA prioritization.'
        }
    ],
    logo_path
)

# ====== SLIDE 6: RULES ENGINE & COMPLIANCE ======
add_content_slide(
    "Detection Rules & Compliance",
    [
        {
            'heading': '⚙️ Customizable Rules Engine',
            'content': 'Rule types: structuring (split transactions), dormant activation, high-risk geography, PEP transactions, velocity anomalies, card fraud, unusual channels. JSON-based conditions support AND/OR logic and 10 evaluation operators. Enable/disable rules by bank policy.'
        },
        {
            'heading': '📋 Regulatory Compliance',
            'content': 'Full compliance with RBI regulatory directives. Auto-generate CTR/SAR reports with mandatory fields. Immutable audit trail proves regulatory requirements met. 5-year data retention with secure encryption and tamper protection.'
        },
        {
            'heading': '🔐 Data Protection',
            'content': 'PII masking in APIs (PAN: XXXXX1234, phone: +91XXXXX3210). Field-level encryption for sensitive data. HTTPS + HSTS enforcement. Rate limiting prevents brute force attacks (5 login attempts/min).'
        }
    ],
    logo_path
)

# ====== SLIDE 7: ADVANCED FEATURES ======
add_content_slide(
    "Advanced Capabilities",
    [
        {
            'heading': '🤖 Investigation Copilot',
            'content': 'AI-powered assistant for investigators. Analyzes alert context, customer history, and rules that triggered. Suggests next investigative steps with confidence scores. Reduces investigation time from hours to minutes.'
        },
        {
            'heading': '🕸️ Network Analysis',
            'content': 'Graph visualization of customer relationships and transaction flows. Identify layering rings and money mule networks. Detect collusion patterns across accounts and customers. Risk concentration analysis.'
        },
        {
            'heading': '📊 Risk Management',
            'content': 'Risk appetite heatmaps by geography, merchant, customer segment. SLA compliance dashboard: target vs actual. Regulatory filing deadlines with countdown. Production-ready monitoring with Prometheus metrics and alerting.'
        }
    ],
    logo_path
)

# ====== SLIDE 8: RISK SCORE CALCULATION ======
add_content_slide(
    "Risk Score Calculation Model",
    [
        {
            'heading': '💹 Transaction Risk Score (0-100)',
            'content': 'Base: 10pt | Rule Severity: Critical +40, High +25, Medium +15, Low +5 | Channel: SWIFT +10, Branch/ATM +5, Internet/Mobile +3, POS +2 | Customer Category: very_high +35, high +20, medium +10, low 0. Example: 10+25+10+35=80/100'
        },
        {
            'heading': '👥 Customer Risk Score (Weighted Rolling)',
            'content': 'Previous score + adjustment per matched rule (capped at 100). Critical rule: +5, High: +3, Medium: +1.5, Low: +0.5. Example: Score 70 + two critical rules = min(70+5+5, 100) = 80/100'
        },
        {
            'heading': '🎯 Auto-Assigned Risk Categories',
            'content': '75-100: very_high (🔴 strict monitoring) | 50-74: high (🟠 enhanced rules) | 25-49: medium (🟡 standard) | 0-24: low (🟢 minimal). Higher scores = stricter alert sensitivity, automatic.'
        }
    ],
    logo_path
)

# ====== SLIDE 9: RISK APPETITE METRICS ======
add_content_slide(
    "Risk Appetite Monitoring — CRO Dashboard",
    [
        {
            'heading': '📊 5 Portfolio-Level Metrics',
            'content': '1️⃣ High-Risk Customer Exposure: (very_high / total) × 100% [Limit: 15% | Warning: 12%] | 2️⃣ Portfolio Avg Risk: Mean risk_score [Limit: 45 | Warning: 35] | 3️⃣ PEP Exposure: (PEP / total) × 100% [Limit: 5% | Warning: 3%]'
        },
        {
            'heading': '🎯 SLA & Transaction Metrics',
            'content': '4️⃣ SLA Compliance Rate: (on-time / total) × 100% [Limit: 85% | Warning: 90%] | 5️⃣ Flagged Txn Volume (30d): (flagged amount / total) × 100% [Limit: 2% | Warning: 1.5%]'
        },
        {
            'heading': '⚠️ Status Indicators & CRO Control',
            'content': '🟢 OK (within warning) | 🟡 WARNING (warning–limit) | 🔴 BREACH (exceeds limit). CRO adjusts thresholds on-the-fly via dashboard—no code changes, instant effect.'
        }
    ],
    logo_path
)

# Save presentation
output_path = r"d:\Claude Projects\EnterpriseRiskSystem\FinsurgeENRIMS_Features_v2.pptx"
prs.save(output_path)
print(f"Created: {output_path}")
print(f"9 slides with risk score and risk appetite metrics")
