from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import MSO_ANCHOR
from pptx.dml.color import RGBColor

# Finsurge brand colors
NAVY_BLUE = RGBColor(31, 67, 128)      # #1F4380
PURPLE = RGBColor(163, 31, 104)        # #A31F68
ACCENT_PINK = RGBColor(234, 61, 117)   # #EA3D75
WHITE = RGBColor(255, 255, 255)
LIGHT_GRAY = RGBColor(245, 246, 248)   # #F5F6F8
DARK_GRAY = RGBColor(55, 65, 81)       # #374151

def add_logo(slide, logo_path):
    """Add Finsurge logo to top-right"""
    left = Inches(8.8)
    top = Inches(0.3)
    height = Inches(0.6)
    slide.shapes.add_picture(logo_path, left, top, height=height)

def add_title_slide(prs, title, subtitle, logo_path):
    """Professional title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background: Navy blue
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = NAVY_BLUE

    # Add Finsurge logo
    add_logo(slide, logo_path)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.5), Inches(8), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.8), Inches(4.2), Inches(8), Inches(2))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.word_wrap = True
    p = subtitle_frame.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(24)
    p.font.color.rgb = ACCENT_PINK

    # Tagline
    tagline_box = slide.shapes.add_textbox(Inches(0.8), Inches(6.5), Inches(8), Inches(0.8))
    tagline_frame = tagline_box.text_frame
    p = tagline_frame.paragraphs[0]
    p.text = "Enterprise Risk Management System for Indian Banking"
    p.font.size = Pt(16)
    p.font.italic = True
    p.font.color.rgb = RGBColor(200, 200, 200)

def add_content_slide(prs, title, content_items, logo_path):
    """Professional content slide with McKinsey styling"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # White background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE

    # Top bar (navy blue accent)
    top_bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.08))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = NAVY_BLUE
    top_bar.line.color.rgb = NAVY_BLUE

    # Add logo
    add_logo(slide, logo_path)

    # Title with accent line
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.35), Inches(8.2), Inches(0.95))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    title_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(38)
    p.font.bold = True
    p.font.color.rgb = NAVY_BLUE

    # Accent line under title
    line = slide.shapes.add_connector(1, Inches(0.8), Inches(1.4), Inches(3.5), Inches(1.4))
    line.line.color.rgb = ACCENT_PINK
    line.line.width = Pt(4)

    # Content box
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.75), Inches(8.4), Inches(5.25))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    text_frame.vertical_anchor = MSO_ANCHOR.TOP

    for i, item in enumerate(content_items):
        if i > 0:
            text_frame.add_paragraph()
        p = text_frame.paragraphs[i]

        # Check if it's a sub-item (starts with 4 spaces)
        if item.startswith("    "):
            p.text = item.strip()
            p.font.size = Pt(16)
            p.level = 1
        else:
            p.text = item
            p.font.size = Pt(18)
            p.level = 0

        p.font.color.rgb = DARK_GRAY
        p.space_before = Pt(6)
        p.space_after = Pt(6)

        # Bold first part before dash for sub-items
        if " — " in p.text:
            runs = p.runs
            if runs:
                parts = p.text.split(" — ")
                if len(parts) == 2:
                    p.text = ""
                    run1 = p.add_run()
                    run1.text = parts[0]
                    run1.font.bold = True
                    run2 = p.add_run()
                    run2.text = " — " + parts[1]

def add_two_column_slide(prs, title, left_title, left_items, right_title, right_items, logo_path):
    """Two-column layout slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE

    top_bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.08))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = NAVY_BLUE
    top_bar.line.color.rgb = NAVY_BLUE

    add_logo(slide, logo_path)

    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.35), Inches(8.2), Inches(0.95))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    title_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(38)
    p.font.bold = True
    p.font.color.rgb = NAVY_BLUE

    line = slide.shapes.add_connector(1, Inches(0.8), Inches(1.4), Inches(3.5), Inches(1.4))
    line.line.color.rgb = ACCENT_PINK
    line.line.width = Pt(4)

    # Left column
    left_header = slide.shapes.add_textbox(Inches(0.8), Inches(1.75), Inches(4), Inches(0.35))
    left_h = left_header.text_frame
    p = left_h.paragraphs[0]
    p.text = left_title
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = PURPLE

    left_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.15), Inches(4.2), Inches(4.85))
    left_frame = left_box.text_frame
    left_frame.word_wrap = True
    for i, item in enumerate(left_items):
        if i > 0:
            left_frame.add_paragraph()
        p = left_frame.paragraphs[i]
        p.text = item
        p.font.size = Pt(15)
        p.font.color.rgb = DARK_GRAY
        p.space_before = Pt(4)
        p.space_after = Pt(4)

    # Right column
    right_header = slide.shapes.add_textbox(Inches(5.2), Inches(1.75), Inches(4), Inches(0.35))
    right_h = right_header.text_frame
    p = right_h.paragraphs[0]
    p.text = right_title
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = PURPLE

    right_box = slide.shapes.add_textbox(Inches(5.2), Inches(2.15), Inches(4.0), Inches(4.85))
    right_frame = right_box.text_frame
    right_frame.word_wrap = True
    for i, item in enumerate(right_items):
        if i > 0:
            right_frame.add_paragraph()
        p = right_frame.paragraphs[i]
        p.text = item
        p.font.size = Pt(15)
        p.font.color.rgb = DARK_GRAY
        p.space_before = Pt(4)
        p.space_after = Pt(4)

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

logo_path = r"d:\Users\lsgop\Downloads\finsurgelogo.png"

# Slide 1: Title
add_title_slide(prs, "FinsurgeENRIMS", "Real-Time Fraud Detection & Risk Management", logo_path)

# Slide 2: System Overview
add_content_slide(prs, "System Overview", [
    "Real-time detection of fraud, cyber attacks, and regulatory violations for Indian banks",
    "77 detection rules organized into 17 detection scenarios covering:",
    "    Fraud patterns (card fraud, account takeover, behavioral anomalies, mule accounts)",
    "    Cyber threats (SIM swap, phishing, UPI fraud, digital channel attacks)",
    "    Advanced fraud (AI deepfakes, synthetic identities, bot detection)",
    "    Internal fraud (employee abuse, data theft)",
    "Customer 360° unified view with network relationship analysis and risk screening",
    "Alert-to-case-to-investigation workflow with SLA tracking and AI-powered investigation copilot"
], logo_path)

# Slide 3: Investigation Copilot
add_content_slide(prs, "Investigation Copilot", [
    "AI-assisted investigation platform that accelerates alert analysis and case resolution",
    "Score Calculation Transparency",
    "    Hover tooltip shows exact risk score formula: Base Alert + Risk Factors + PEP Flag + Severity",
    "    Helps investigators understand why alerts are prioritized",
    "Ask the Copilot — Conversational Analysis",
    "    'Why is this customer flagged?' — Copilot explains detection pattern and rule triggers",
    "    'Show similar cases' — Retrieves historical closed cases with similar characteristics",
    "    'Next steps?' — Recommends investigation actions based on alert profile",
    "Context Integration: Previous cases, activity timeline, related alerts all visible in one view"
], logo_path)

# Slide 4: Detection Scenarios
add_content_slide(prs, "Detection Scenarios — Read-Only Framework", [
    "Detection Scenarios are read-only templates that group related rules by fraud/risk type",
    "17 Pre-Built Scenarios (not editable via UI):",
    "    Fraud: Card Fraud (3 rules), Account Takeover (3), Wire Fraud (3)",
    "    Cyber: SIM Swap/Phishing (3), Digital Channel Attacks (4), UPI/Instant Payment (3)",
    "    Advanced: Deepfake/Synthetic ID (5), Bot Detection (3)",
    "    Internal: Employee Abuse (7 rules), Data Theft (3)",
    "    Compliance: Dormant Account (1), Beneficiary Risk (1)",
    "Scenario Metrics: Detection count, last triggered date, constituent rules status",
    "Purpose: Provide investigators immediate context about fraud pattern types without manual setup"
], logo_path)

# Slide 5: Rules Engine
add_content_slide(prs, "Rules Engine — Transaction Evaluation", [
    "Real-time evaluation of 77 enabled detection rules against every customer transaction",
    "Rule Architecture",
    "    Conditions: JSON-based logic trees (e.g., IF age > 30 AND amount > 10L, THEN flag)",
    "    Actions: Create alert, flag transaction, adjust customer risk score, notify compliance",
    "    Thresholds: Amount-based, count-based, time-windowed (1h, 24h, 7d, 30d, 90d)",
    "    Channels: Apply to specific channels (Mobile, Teller, ATM, Card, UPI)",
    "    Customer Types: Target rule to specific segments (High-Value, New, Dormant)",
    "On-Rule Match Workflow",
    "    Alert created with risk score, assigned to analyst based on workload, SLA timer starts",
    "Performance Tracking: Detection count per rule, false positive rate, rule effectiveness"
], logo_path)

# Slide 6: Customer 360 & Network
add_content_slide(prs, "Customer 360 — Complete Risk Profile", [
    "Single-pane view of customer risk, transactions, cases, and connected entities",
    "5 Integrated Tabs:",
    "    Overview — Risk score, KYC status, account summary, transaction heatmap",
    "    Transactions — Full transaction history with rule-match filtering",
    "    Alerts — All linked alerts with status, SLA, and investigation notes",
    "    Network Analysis — Relationship mapping showing connected entities",
    "    Audit — Immutable log of all customer-related actions and data access",
    "Network Relationship Types: Parent/Subsidiary, UBO (Ultimate Beneficial Owner), Director, Shareholder",
    "Risk Visualizations: Entity risk scores, PEP flags, sanctions match indicators",
    "Enables investigators to understand customer's business ecosystem in seconds"
], logo_path)

# Slide 7: Alert Management & Investigation
add_content_slide(prs, "Alert Management — From Detection to Action", [
    "Complete alert lifecycle with SLA enforcement and intelligent routing",
    "Alert States: New — Assigned — Under Review — Escalated — Closed (True/False Positive)",
    "Severity-Based SLAs: Critical (4h), High (24h), Medium (72h), Low (168h response time)",
    "Intelligent Auto-Assignment: Routes alerts to available analysts by skill, current workload",
    "Investigation Features",
    "    Add detailed investigation notes, attach evidence documents, link related cases",
    "    Bulk actions: Close multiple false positives, escalate groups, reassign batches",
    "    Override alerts, mark as false positive with reason, escalate for compliance review",
    "Dashboard Analytics: New alerts trending, SLA compliance rate, closure patterns by rule/analyst"
], logo_path)

# Slide 8: Case Management & Regulatory Filing
add_content_slide(prs, "Case Management & Document Downloads", [
    "Consolidate related alerts and evidence into formal investigation cases",
    "Case Types: Fraud, Cyber Attack, Money Laundering, Operational Risk, Internal Misconduct",
    "Case Lifecycle: Open — Assigned — Under Investigation — Escalated — Pending Resolution — Closed",
    "Document Generation & Download",
    "    Police FIR — Download full investigation report with all case details",
    "    CTR/SAR Filings — Download regulatory reports for archival or reprint",
    "    Generate formatted .txt documents suitable for filing and compliance",
    "Evidence Management: Attach documents, screenshots, transaction exports, supporting records",
    "Integration: Auto-escalate when case confirmed as fraud or suspicious activity"
], logo_path)

# Slide 9: Regulatory Filing & Dashboard
add_content_slide(prs, "Regulatory Filing & Real-Time Monitoring", [
    "Automated filing workflow with document generation and download capabilities",
    "CTR/SAR Filings: Download reports with complete filing details and metadata",
    "FIR Management: File reports with police, track investigation progress, download documents",
    "Executive visibility into risk management operations and team performance",
    "Key Performance Indicators (Real-Time)",
    "    New alerts (24h), Overdue alerts, New cases, Closure rate, SLA compliance %",
    "Operational Dashboards",
    "    Alert heatmap by hour, customer segment, rule category showing risk distribution",
    "    Team workload: Assignments per analyst, case load, average case resolution time",
    "    Rules performance: Top-triggered rules, false positive rate, detection effectiveness"
], logo_path)

# Slide 10: Audit Trail & Security
add_content_slide(prs, "Audit Trail — Compliance & Accountability", [
    "Immutable audit logging of every operation for regulatory compliance",
    "Audit Coverage",
    "    Detection: Rule creation/modification, alert generation, escalation triggers",
    "    Investigation: Case creation, evidence attachment, notes added, disposition marked",
    "    Access: User login, data access, rule changes, configuration updates",
    "    Document Downloads: Track all filed document and report downloads for compliance",
    "    Retention: 5 years (RBI mandate) for audit logs, 7 years for transaction data",
    "Security Controls",
    "    Role-Based Access: Admin, Compliance, Analyst, Investigator, Maker-Checker approval",
    "    Data Protection: PII encryption (PAN, Aadhaar), API response masking, HTTPS-only",
    "    Compliance Ready: Meets RBI Master Direction, PMLA Section 12, data localization requirements"
], logo_path)

# Save
prs.save("FinsurgeENRIMS_Professional_Updated.pptx")
print("[OK] Updated PowerPoint created: FinsurgeENRIMS_Professional_Updated.pptx")
print("[FEATURE] Added document download capabilities for CTR, SAR, and FIR")
print("[FEATURE] Enhanced Case Management slide to showcase document downloads")
print("[FEATURE] Updated Regulatory Filing slide with download documentation")
print("[UPDATE] 10 slides with latest system features")
