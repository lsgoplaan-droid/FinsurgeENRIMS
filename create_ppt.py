from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def add_title_slide(prs, title, subtitle):
    """Add a title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(31, 41, 55)

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(9), Inches(2))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.word_wrap = True
    p = subtitle_frame.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(28)
    p.font.color.rgb = RGBColor(209, 213, 219)

def add_content_slide(prs, title, content_list):
    """Add a content slide with bullet points"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(31, 41, 55)

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(9), Inches(5.5))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True

    for i, item in enumerate(content_list):
        if i > 0:
            text_frame.add_paragraph()
        p = text_frame.paragraphs[i]
        p.text = item
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(229, 231, 235)
        p.space_before = Pt(8)
        p.space_after = Pt(8)
        p.level = 0

prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Slide 1: Title
add_title_slide(prs, "FinsurgeENRIMS", "Enterprise Risk Management System\nCore Functions & Features")

# Slide 2: System Overview
add_content_slide(prs, "System Overview", [
    "Real-time fraud detection, AML compliance, and KYC management for Indian banks",
    "92 detection rules organized into 22 scenarios (structuring, card fraud, cyber, AI fraud, etc.)",
    "Customer 360-degree view with network analysis, watchlist screening, risk scoring",
    "Alert-driven investigation workflow with AI-powered copilot assistance",
    "Case management, regulatory reporting (CTR/SAR), immutable audit trails",
    "Role-based access control (Admin, Compliance, Analyst, Investigator, Maker-Checker)"
])

# Slide 3: Investigation Copilot
add_content_slide(prs, "Investigation Copilot - AI Assistant", [
    "AI-powered assistant helping analysts investigate suspicious alerts in real-time",
    "Score Calculation Tooltip - Shows how alert risk score is computed:",
    "    • Base alert risk + Risk factors + PEP status + Severity = Final score",
    "Ask the Copilot - Natural language chat to analyze alerts:",
    "    • Why is this customer flagged? Copilot explains the pattern",
    "    • Similar cases? Fetches related closed cases from history",
    "Displays previous cases, activity timeline, recommended actions",
    "Live score breakdown, contextual recommendations, pattern matching"
])

# Slide 4: Detection Scenarios & Rules Engine
add_content_slide(prs, "Detection Scenarios & Rules Engine", [
    "Detection Scenarios - Named rule groupings for specific fraud patterns:",
    "    AML: Structuring (5), Layering (4), Geographic Risk (3)",
    "    Fraud: Card Fraud (3), Account Takeover (3), Behavioral (5), Mule (3)",
    "    Cyber: SIM Swap/Phishing (3), Digital Attacks (4), UPI Fraud (3)",
    "    Internal: Employee Abuse (7), Data Theft (3) | AI: Deepfake (5), Bot (3)",
    "Rules Engine - Evaluates 75 enabled rules against each transaction",
    "    JSON condition trees, 10+ operators, AND/OR nesting",
    "    On match: Creates alert, flags transaction, adjusts risk, auto-assigns analyst"
])

# Slide 5: Customer 360 - Complete Profile
add_content_slide(prs, "Customer 360 - Complete Customer Profile", [
    "Unified customer view with 5 tabs: Overview, Transactions, Alerts, Network, Audit",
    "Overview - Risk score, KYC status, account summary, activity heatmap",
    "Network Analysis Tab - Connected entities with relationship types:",
    "    Parent companies, subsidiaries, UBOs, directors, shareholders",
    "    Risk badges: High-risk entities, PEPs, sanctions matches",
    "    Network summary: Total entities, connections, average risk score",
    "Transactions & Alerts - Drill-down into specific activity",
    "Audit Tab - Immutable log of all customer-related actions"
])

# Slide 6: Alert Management - Detect & Investigate
add_content_slide(prs, "Alert Management - Detect & Investigate", [
    "Real-time alerts when rules trigger on customer transactions",
    "Alert States: New > Assigned > Under Review > Escalated > Closed (True/False Positive)",
    "SLA Tracking: Critical (4h), High (24h), Medium (72h), Low (168h)",
    "Auto-assignment: New alerts routed to analysts by skill & workload",
    "Alert Details: Risk score, triggered rule, customer context, recommendations",
    "Investigation: Add notes, attach evidence, mark disposition, escalate if needed",
    "Bulk Operations: Close false positives, escalate groups, re-assign batches",
    "Dashboard: Real-time KPIs (new alerts, overdue, closed rate)"
])

# Slide 7: Case Management - Investigation
add_content_slide(prs, "Case Management - Structured Investigation", [
    "Cases - Consolidate multiple related alerts into formal investigations",
    "Links: Multiple alerts, customers, evidence documents, CTR/SAR reports",
    "Workflow: Open > Assigned > Under Investigation > Escalated > Pending Regulatory > Closed",
    "Case Types: Fraud, AML/Structuring, KYC, Operational, Cyber, Internal Fraud",
    "Evidence Management: Attach documents, screenshots, transaction exports",
    "Regulatory Tie-in: Auto-file CTR or SAR based on investigation findings",
    "Metrics: Status distribution, investigation time, true positive rate",
    "Full Audit: Creation, reassignments, status changes all logged immutably"
])

# Slide 8: Rules Management - Customize Detection
add_content_slide(prs, "Rules Management - Customize Detection", [
    "Rule Editor: Create, modify, enable/disable rules without code",
    "Rule Components: Conditions (JSON trees), Actions, Time Window, Thresholds",
    "    Example: If customer_age > 30 AND txn_amount > 10L, create high-priority alert",
    "Time Windows: 1h, 24h, 7d, 30d, 90d aggregation periods",
    "Thresholds: Amount (paise), count, applicable channels, customer types",
    "Metrics View: Detection count, last triggered, false positive rate per rule",
    "75 Enabled Rules: Fraud, Cyber, Internal, AI Fraud categories",
    "Versioning: Track changes, rollback capability, maker-checker approval required"
])

# Slide 9: KYC & Compliance - Regulatory
add_content_slide(prs, "KYC & Compliance - Regulatory Adherence", [
    "KYC Reviews: Initial, annual re-KYC, periodic reviews with documents",
    "Documents: Aadhaar, PAN, proof of address, video KYC status verification",
    "Risk Classification: Based on profile, PEP check, sanctions screening, patterns",
    "Watchlist Screening: Real-time match against PEPs, sanctions, internal lists",
    "Screening Results: Match status, confidence scores, exemption marking",
    "Audit Trail: Every KYC update, screening run, exemption with user, timestamp, reason",
    "Data Protection: KYC data localized to India, encrypted at rest, 7-year retention",
    "RBI Compliance: Meets Master Direction on KYC, AML/CFT standards"
])

# Slide 10: Watchlist & Reporting - CTR/SAR Filing
add_content_slide(prs, "Watchlist & Regulatory Reporting", [
    "Watchlist Management: Internal lists of high-risk entities, PEPs, blocked individuals",
    "Screening: Every customer auto-screened; matches flag for review",
    "Watchlist Types: PEPs, Sanctions, Debarred, High-Risk Jurisdictions, Corporate",
    "CTR Filing: Cash Transaction Reports auto-generated for txns >= 10 lakh",
    "SAR Filing: Suspicious Activity Reports filed for confirmed fraud/AML cases",
    "Filing Status: Track submissions to FIU-IND with reference numbers, dates",
    "Compliance: Report counts by period, filing timelines, rejections/resubmissions",
    "Integration: Direct FIU-IND portal integration for automated electronic filing"
])

# Slide 11: Dashboard & Monitoring
add_content_slide(prs, "Dashboard & Monitoring - Real-time Visibility", [
    "KPI Cards: Active alerts, overdue, new cases, closed rate, SLA performance",
    "Alert Heatmap: Risk distribution by hour, customer segment, rule category",
    "Rules Performance: Top-triggered rules, false positive rate, effectiveness",
    "Team Workload: Alerts per analyst, case load, average closure time",
    "Audit Log Summary: Recent activities, access patterns, failed auth attempts",
    "System Health: Database status, API latency, background job status",
    "Geographic Heat: Transaction volume and risk concentration by branch/state",
    "Customizable Dashboards: Saved views by role, filters, date ranges"
])

# Slide 12: Audit Trail & Security
add_content_slide(prs, "Audit Trail & Security - Compliance", [
    "Immutable Audit Log: Every operation logged (user, IP, timestamp, resource, action)",
    "Audit Coverage: Rule changes, user roles, alert closures, case updates, data access",
    "Role-Based Access: Admin, Compliance, Analyst, Investigator, Maker-Checker",
    "Data Security: PII encryption (PAN, Aadhaar), response masking, HTTPS-only",
    "Retention: Audit logs 5 years (RBI), transaction data 7 years (RBI mandate)",
    "Audit Reports: Export by date, user, action type for compliance reviews",
    "Compliance Ready: RBI Master Direction, PMLA Section 12, data localization",
    "Production Standards: Meets security audit requirements for bank deployment"
])

prs.save("FinsurgeENRIMS_Functions_Guide.pptx")
print("PowerPoint created: FinsurgeENRIMS_Functions_Guide.pptx (12 slides)")
