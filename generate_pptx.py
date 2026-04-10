from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
import os

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Finsurge colors
PRIMARY_BLUE = RGBColor(75, 85, 180)
ACCENT_ORANGE = RGBColor(255, 102, 0)
DARK_TEXT = RGBColor(40, 40, 40)
LIGHT_BG = RGBColor(245, 247, 250)
WHITE = RGBColor(255, 255, 255)

def add_title_slide(title, subtitle, logo_path=None):
    """Add title slide with logo"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE

    # Add gradient-like background with colored bar
    shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = PRIMARY_BLUE
    shape.line.color.rgb = PRIMARY_BLUE

    # Add logo if exists
    if logo_path and os.path.exists(logo_path):
        try:
            slide.shapes.add_picture(logo_path, Inches(0.5), Inches(0.5), width=Inches(1.5))
        except:
            pass

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(54)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = WHITE

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(9), Inches(3))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.word_wrap = True
    subtitle_frame.text = subtitle
    subtitle_frame.paragraphs[0].font.size = Pt(28)
    subtitle_frame.paragraphs[0].font.color.rgb = DARK_TEXT

def add_content_slide(title, sections, logo_path=None, notes=""):
    """Add content slide with multiple sections"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = LIGHT_BG

    # Add logo in footer (bottom right) if exists
    if logo_path and os.path.exists(logo_path):
        try:
            slide.shapes.add_picture(logo_path, Inches(8.8), Inches(6.8), width=Inches(1))
        except:
            pass

    # Add speaker notes if provided
    if notes:
        text_frame = slide.notes_slide.notes_text_frame
        text_frame.text = notes

    # Title bar
    title_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.8))
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = PRIMARY_BLUE
    title_shape.line.color.rgb = PRIMARY_BLUE

    # Title text
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(8), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(40)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = WHITE

    # Add sections
    y_position = 1.2
    for i, section in enumerate(sections):
        # Section heading with orange accent
        accent = slide.shapes.add_shape(1, Inches(0.3), Inches(y_position), Inches(0.1), Inches(0.4))
        accent.fill.solid()
        accent.fill.fore_color.rgb = ACCENT_ORANGE
        accent.line.color.rgb = ACCENT_ORANGE

        heading_box = slide.shapes.add_textbox(Inches(0.6), Inches(y_position), Inches(9), Inches(0.35))
        heading_frame = heading_box.text_frame
        heading_frame.text = section['heading']
        heading_frame.paragraphs[0].font.size = Pt(18)
        heading_frame.paragraphs[0].font.bold = True
        heading_frame.paragraphs[0].font.color.rgb = PRIMARY_BLUE

        # Section content
        content_box = slide.shapes.add_textbox(Inches(0.6), Inches(y_position + 0.4), Inches(9), Inches(1.2))
        content_frame = content_box.text_frame
        content_frame.word_wrap = True
        content_frame.text = section['content']
        for para in content_frame.paragraphs:
            para.font.size = Pt(14)
            para.font.color.rgb = DARK_TEXT
            para.space_before = Pt(3)
            para.space_after = Pt(3)

        y_position += 1.6
        if y_position > 6.5:
            break

logo_path = r"D:\Users\lsgop\Downloads\finsurgelogo.png"

# ====== SLIDE 1: TITLE SLIDE ======
add_title_slide(
    "FinsurgeENRIMS",
    "Enterprise Risk Management System\nReal-time Detection, Investigation & Regulatory Compliance for Banking",
    logo_path
)

# ====== SLIDE 2: DASHBOARD ======
dashboard_notes = """Dashboard Details:
- KPI tiles show real-time metrics updated every 5 minutes
- Open alerts: All non-closed statuses (new, assigned, escalated, under_investigation)
- SLA response times: Critical (2h), High (4h), Medium (8h), Low (24h)
- Fraud loss calculated as sum of transaction amounts flagged as suspicious
- All metrics are clickable drill-downs to detailed list views with full filtering
- Trend charts show 30-day, 90-day, and year-to-date views"""
add_content_slide(
    "Dashboard — Real-Time Oversight",
    [
        {
            'heading': '📊 Executive KPIs',
            'content': 'Centralized view of critical metrics: open alerts (30), active cases (12), fraud loss (₹53.5M), overdue SLA (2 cases). Color-coded severity indicators for quick identification of risk areas.'
        },
        {
            'heading': '📈 Trend Analysis',
            'content': 'Visualize alert volume, case closure rates, and regulatory filing trends over time. Drill-down capability to trace metrics back to source transactions and customers.'
        },
        {
            'heading': '⚡ Instant Actions',
            'content': 'Click any metric to filter and investigate. Create cases directly from alerts. Assign alerts to analysts with SLA tracking (response time: 2-24 hours based on priority).'
        }
    ],
    logo_path,
    dashboard_notes
)

# ====== SLIDE 3: ALERT MANAGEMENT ======
alert_notes = """Alert Management Details:
- 26 customizable rules: structuring, dormant activation, PEP transactions, velocity, geography, card fraud, channel anomalies
- Each rule has configurable thresholds (amounts, time windows, frequencies)
- New alerts assigned automatically by Investigation Copilot AI to available analysts
- Alert lifecycle: New → Under Investigation → Escalated or Closed → Case Created or Dismissed
- Response SLA enforced: Critical alerts 2h, High 4h, Medium 8h, Low 24h
- Full audit trail maintained for regulatory review (5-year retention)
- Bulk operations supported: bulk close, bulk escalate with audit justification"""
add_content_slide(
    "Alert Detection & Investigation",
    [
        {
            'heading': '🎯 Automated Detection',
            'content': '26 detection rules monitor transactions 24/7 for suspicious patterns: structuring, dormant account activation, unusual velocity, high-risk geography. Real-time risk scoring (0-100) on every transaction.'
        },
        {
            'heading': '🔍 Investigation Copilot',
            'content': 'AI-powered analysis of new alerts. Automatic context gathering: customer profile, account history, network relationships. Smart assignment to available analysts with workload balancing.'
        },
        {
            'heading': '✅ Alert Response',
            'content': 'Mark as true positive → escalate to cases. False positive → close with reason. Escalated → manager review. Full audit trail with user, timestamp, and notes on every status change.'
        }
    ],
    logo_path,
    alert_notes
)

# ====== SLIDE 4: CASE MANAGEMENT ======
case_notes = """Case Management Workflow:
- Case statuses: New → Assigned → Under Investigation → Escalated → Closed (True Positive, False Positive, Inconclusive)
- Disposition options: True Positive (suspicious confirmed), False Positive (no risk), Inconclusive (insufficient data)
- SLA tracking: Open SLA starts on creation, escalation resets SLA to manager's availability
- Evidence attachment: Supports documents (PDF), screenshots, bank statements, investigator comments
- Case linking: Multiple alerts can be linked to single case (e.g., related structuring transactions)
- Timeline: Auto-generated from all case activities with timestamps and user attribution
- True Positive cases auto-generate CTR/SAR with pre-filled transaction details
- Mandatory fields in CTR/SAR: Transaction amount, customer info, reason for suspicion, investigation findings"""
add_content_slide(
    "Case Management & Escalation",
    [
        {
            'heading': '📋 Full Case Lifecycle',
            'content': 'Create case from alert → assign to investigator → track findings → escalate if needed → set disposition (true/false positive) → file regulatory report. SLA tracking with visual indicators for overdue cases.'
        },
        {
            'heading': '🔗 Evidence & Context',
            'content': 'Attach evidence (documents, screenshots, investigator notes). Link related alerts and transactions. Build investigation narrative with automated timeline from case activity log.'
        },
        {
            'heading': '⚖️ Regulatory Workflow',
            'content': 'Cases flagged as true positive auto-file CTR/SAR to FIU-IND portal. Track filing status and regulatory feedback. Maintain 5-year audit trail for compliance reviews.'
        }
    ],
    logo_path,
    case_notes
)

# ====== SLIDE 5: CUSTOMER 360 & FRAUD DETECTION ======
risk_scoring_notes = """DETAILED RISK SCORING METHODOLOGY:

BASE RISK CATEGORIES (Starting Points):
- Low (0-25): Individual customers with clean transaction history, verified employment
- Medium (26-50): Newly onboarded customers or those with occasional suspicious activity
- High (51-75): PEP customers, high-transaction volume, emerging markets, cash-heavy businesses
- Very High (76-100): Multiple red flags, regulatory warnings, known high-risk jurisdictions

BEHAVIORAL MODIFIERS (Applied Dynamically):
1. Transaction Velocity: Amount and frequency patterns
   - Spike detection: +15 points if current volume > 3x rolling 30-day average
   - Dormancy reversal: +20 points if account inactive >90 days then large transaction

2. Geographic Indicators: +10-30 points for high-risk jurisdictions
   - Sanctioned countries, high-corruption index, FIU watchlist

3. Network Risk: +5-25 points based on connected accounts
   - Unusual beneficiary relationships, circular transfers

4. Merchant Category: +10-20 points for high-risk sectors
   - Money services, gems/precious metals, casinos, import-export

FINAL SCORE CALCULATION:
Risk Score = Base Risk + Σ(Behavioral Modifiers) with caps at 0-100

ALERT SENSITIVITY IMPACT:
- Low risk (0-25): Rules trigger at strict thresholds only (high false positive cost)
- Medium (26-50): Standard thresholds, balanced sensitivity
- High (51-75): Enhanced sensitivity, lower thresholds trigger alerts
- Very High (76-100): Maximum sensitivity, aggressive monitoring, mandatory escalation

SLA IMPACT:
- Low/Medium: 24-hour response SLA
- High: 8-hour response SLA
- Very High: 2-hour response SLA, automatic escalation to senior analyst

SCORE RECALCULATION:
- Triggered on every transaction (real-time)
- Weekly review of behavioral modifiers
- Manual adjustment by compliance officer with full audit trail"""
add_content_slide(
    "Customer 360 & Fraud Detection",
    [
        {
            'heading': '👤 Customer 360 Profile',
            'content': 'Complete customer view: verification status, risk category, accounts, transaction history (last 50), related alerts, open cases, network relationships. PEP (politically exposed person) flagging with auto-enhanced monitoring for high-risk profiles.'
        },
        {
            'heading': '💳 Fraud Detection',
            'content': 'Card fraud monitoring: velocity checks (amount/frequency), cross-border thresholds, merchant category anomalies. Transaction scoring integrates customer risk + behavioral patterns + rules engine output.'
        },
        {
            'heading': '📊 Risk Scoring',
            'content': 'Dynamic risk adjustment: base risk (low/medium/high/very_high) + behavioral modifiers = real-time customer risk score (0-100). Influences alert sensitivity and SLA prioritization.'
        }
    ],
    logo_path,
    risk_scoring_notes
)

# ====== SLIDE 6: RULES ENGINE & COMPLIANCE ======
rules_notes = """Rules Engine & Compliance Details:
- 26 default rules covering RBI Master Direction requirements
- Rule conditions: JSON syntax with operators (>, <, ==, IN, CONTAINS, AND, OR, NOT, BETWEEN, CONTAINS_ALL)
- Each rule has: trigger threshold, time window (24h, 7d, 30d), customer/account scope
- Rule evaluation happens in real-time on every transaction (< 100ms latency)
- Rules can be enabled/disabled per bank policy without code changes
- No-code UI for rule management: rule builder, testing, audit trail

Compliance Features:
- CTR (Currency Transaction Report) auto-filed when amount >= ₹10 lakhs
- SAR (Suspicious Activity Report) auto-filed when risk triggers detected
- Audit trail: immutable log of all transactions, alerts, cases, rule changes (5-year retention)
- Data localization: India-only storage (RBI circular 2018 compliant)
- Encryption: PAN/Aadhaar encrypted at rest, PII masked in logs

Security:
- PBKDF2 password hashing (not bcrypt - Python 3.12 compatible)
- JWT tokens with 30-minute expiry + sliding refresh
- Rate limiting: 5 login attempts/min per IP, 100 API requests/min per user
- Field-level encryption for PAN, Aadhaar (searchable encryption)
- HTTPS mandatory in production with HSTS headers (1-year validity)"""
add_content_slide(
    "Detection Rules & Compliance",
    [
        {
            'heading': '⚙️ Customizable Rules Engine',
            'content': 'Rule types: structuring (split transactions), dormant activation, high-risk geography, PEP transactions, velocity anomalies, card fraud, unusual channels. JSON-based conditions support AND/OR logic and 10 evaluation operators. Enable/disable rules by bank policy.'
        },
        {
            'heading': '📋 Regulatory Compliance',
            'content': 'Full compliance with RBI regulatory directives. Auto-generate CTR/SAR reports with mandatory fields. Immutable audit trail proves regulatory requirements met. 5-year data retention with secure encryption and tamper protection.'
        },
        {
            'heading': '🔐 Data Protection',
            'content': 'PII masking in APIs (PAN: XXXXX1234, phone: +91XXXXX3210). Field-level encryption for sensitive data. HTTPS + HSTS enforcement. Rate limiting prevents brute force attacks (5 login attempts/min).'
        }
    ],
    logo_path,
    rules_notes
)

# ====== SLIDE 7: ADVANCED FEATURES ======
advanced_notes = """Advanced Features Details:

Investigation Copilot AI:
- Analyzes: Customer KYC profile, account velocity, prior alerts, transaction patterns, rule triggers
- Provides: Summarized risk profile, key red flags, suggested next steps, confidence scores
- Reduces investigation time: 4-6 hours → 30-45 minutes per alert
- Integration: Available in UI as side panel, no separate tool switching needed

Network Analysis:
- Visualizes: Connected accounts, transaction flows, beneficiary relationships
- Detection: Layering rings (circular transfers), money mule networks, collusion patterns
- Scoring: Risk concentration by customer group, geography hotspots
- Performance: Real-time analysis on 1000s of accounts with millisecond response

Risk Management Dashboard:
- Risk appetite: Configurable thresholds by geography (India/overseas), merchant category, customer type
- Heat maps: Show actual vs target exposure with drill-down to underlying transactions
- SLA monitoring: Real-time tracking of response times vs commitments (2h/4h/8h/24h targets)
- Regulatory timeline: Auto-countdown to filing deadlines for CTR/SAR
- Production readiness: Prometheus metrics for all KPIs, Grafana dashboards, PagerDuty alerting"""
add_content_slide(
    "Advanced Capabilities",
    [
        {
            'heading': '🤖 Investigation Copilot',
            'content': 'AI-powered assistant for investigators. Analyzes alert context, customer history, and rules that triggered. Suggests next investigative steps with confidence scores. Reduces investigation time from hours to minutes.'
        },
        {
            'heading': '🕸️ Network Analysis',
            'content': 'Graph visualization of customer relationships and transaction flows. Identify layering rings and money mule networks. Detect collusion patterns across accounts and customers. Risk concentration analysis.'
        },
        {
            'heading': '📊 Risk Management',
            'content': 'Risk appetite heatmaps by geography, merchant, customer segment. SLA compliance dashboard: target vs actual. Regulatory filing deadlines with countdown. Production-ready monitoring with Prometheus metrics and alerting.'
        }
    ],
    logo_path,
    advanced_notes
)

# Save presentation
output_path = r"d:\Claude Projects\EnterpriseRiskSystem\FinsurgeENRIMS_Features_7Slides.pptx"
prs.save(output_path)
print(f"✅ Created: {output_path}")
print(f"📊 7 slides with professional design and Finsurge branding")
